"""Extract text from uploaded documents and chunk it for the LLM."""
from pathlib import Path

CHUNK_CHARS = 8000        # ~2k tokens per chunk
CHUNK_OVERLAP = 400
MIN_CHUNK_CHARS = 300     # skip near-empty chunks


class ParseError(Exception):
    pass


def extract_text(path: str, filetype: str, max_pages: int) -> list[tuple[str, str]]:
    """Return list of (section_label, text)."""
    p = Path(path)
    if not p.exists():
        raise ParseError("Uploaded file is missing on disk")
    if filetype == "pdf":
        return _pdf(p, max_pages)
    if filetype == "docx":
        return _docx(p)
    if filetype == "pptx":
        return _pptx(p)
    if filetype == "epub":
        return _epub(p)
    if filetype == "md":
        return [("document", p.read_text(encoding="utf-8", errors="replace"))]
    raise ParseError(f"Unsupported file type: {filetype}")


def _pdf(p: Path, max_pages: int) -> list[tuple[str, str]]:
    import fitz  # PyMuPDF
    doc = fitz.open(p)
    if doc.page_count > max_pages:
        raise ParseError(f"PDF has {doc.page_count} pages; max is {max_pages}")
    out = []
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            out.append((f"page {i}", text))
    doc.close()
    if not out:
        raise ParseError("No selectable text found — this PDF may be scanned images (OCR not enabled yet)")
    return out


def _docx(p: Path) -> list[tuple[str, str]]:
    import docx
    d = docx.Document(str(p))
    sections, current_label, buf = [], "start", []
    for para in d.paragraphs:
        if para.style.name.startswith("Heading") and para.text.strip():
            if buf:
                sections.append((current_label, "\n".join(buf)))
                buf = []
            current_label = para.text.strip()[:80]
        elif para.text.strip():
            buf.append(para.text)
    if buf:
        sections.append((current_label, "\n".join(buf)))
    if not sections:
        raise ParseError("No text found in document")
    return sections


def _pptx(p: Path) -> list[tuple[str, str]]:
    from pptx import Presentation
    prs = Presentation(str(p))
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if texts:
            out.append((f"slide {i}", "\n".join(texts)))
    if not out:
        raise ParseError("No text found in presentation")
    return out


def _epub(p: Path) -> list[tuple[str, str]]:
    from bs4 import BeautifulSoup
    from ebooklib import ITEM_DOCUMENT
    from ebooklib import epub

    book = epub.read_epub(str(p))
    documents = {item.get_id(): item for item in book.get_items_of_type(ITEM_DOCUMENT)}
    ordered_ids = [
        item_id
        for item_id, _linear in book.spine
        if isinstance(item_id, str) and item_id in documents
    ]
    if not ordered_ids:
        ordered_ids = list(documents)

    sections = []
    for i, item_id in enumerate(ordered_ids, start=1):
        item = documents[item_id]
        html = item.get_content().decode("utf-8", errors="replace")
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav"]):
            tag.decompose()

        title = None
        heading = soup.find(["h1", "h2", "h3"])
        if heading:
            title = heading.get_text(" ", strip=True)
        if not title:
            title = item.get_name().rsplit("/", 1)[-1].rsplit(".", 1)[0]
        title = title.strip()[:80] or f"chapter {i}"

        text = soup.get_text("\n", strip=True)
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        text = "\n".join(lines)
        if text:
            sections.append((f"chapter {i}: {title}", text))

    if not sections:
        raise ParseError("No text found in EPUB")
    return sections


def chunk_sections(sections: list[tuple[str, str]]) -> list[tuple[str, str]]:
    """Merge/split sections into LLM-sized chunks, keeping a source label."""
    chunks = []
    buf_label, buf = None, ""
    for label, text in sections:
        if len(buf) + len(text) <= CHUNK_CHARS:
            buf_label = buf_label or label
            buf += ("\n\n" if buf else "") + text
        else:
            if buf and len(buf) >= MIN_CHUNK_CHARS:
                chunks.append((f"{buf_label}–{label}", buf))
            # split oversized single sections
            while len(text) > CHUNK_CHARS:
                chunks.append((label, text[:CHUNK_CHARS]))
                text = text[CHUNK_CHARS - CHUNK_OVERLAP:]
            buf_label, buf = label, text
    if buf and len(buf) >= MIN_CHUNK_CHARS:
        chunks.append((buf_label or "document", buf))
    if not chunks:
        raise ParseError("Document too short to generate cards from")
    return chunks
